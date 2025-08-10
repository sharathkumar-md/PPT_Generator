"""
Slide Templates and Themes for RealtyAI PowerPoint Generator
============================================================

Defines themes, layouts, and styling options for PowerPoint presentations.
"""

from typing import Dict, Any, Optional, Tuple
from dataclasses import dataclass
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


@dataclass
class ThemeColors:
    """Color scheme for presentation themes"""
    primary: Tuple[int, int, int] = (68, 114, 196)      # Blue
    secondary: Tuple[int, int, int] = (112, 173, 71)    # Green
    accent: Tuple[int, int, int] = (255, 192, 0)        # Orange
    text_dark: Tuple[int, int, int] = (68, 68, 68)      # Dark Gray
    text_light: Tuple[int, int, int] = (255, 255, 255)  # White
    background: Tuple[int, int, int] = (255, 255, 255)  # White
    
    def get_rgb_color(self, color_name: str) -> RGBColor:
        """Convert color tuple to RGBColor object"""
        color_tuple = getattr(self, color_name, self.primary)
        return RGBColor(*color_tuple)


@dataclass
class ThemeFonts:
    """Font configuration for presentation themes"""
    title_font: str = "Calibri"
    body_font: str = "Calibri"
    title_size: int = 44
    subtitle_size: int = 32
    body_size: int = 18
    bullet_size: int = 16


@dataclass 
class SlideLayout:
    """Configuration for individual slide layouts"""
    layout_id: int
    name: str
    title_placeholder_idx: int = 0
    content_placeholder_idx: int = 1
    
    def __post_init__(self):
        """Validate layout configuration"""
        if self.layout_id < 0:
            raise ValueError("Layout ID must be non-negative")


class Theme:
    """Main theme class that defines the overall presentation appearance"""
    
    def __init__(self, 
                 name: str = "Modern",
                 colors: Optional[ThemeColors] = None,
                 fonts: Optional[ThemeFonts] = None):
        """
        Initialize a presentation theme
        
        Args:
            name: Theme name
            colors: Color scheme (uses default if None)
            fonts: Font configuration (uses default if None)
        """
        self.name = name
        self.colors = colors or ThemeColors()
        self.fonts = fonts or ThemeFonts()
        
        # Define available slide layouts
        self.layouts = {
            'title': SlideLayout(0, "Title Slide", 0, 1),
            'content': SlideLayout(1, "Content with Bullets", 0, 1),
            'section_header': SlideLayout(2, "Section Header", 0, -1),
            'two_content': SlideLayout(3, "Two Content", 0, 1),
            'comparison': SlideLayout(4, "Comparison", 0, 1),
            'blank': SlideLayout(6, "Blank", -1, -1)
        }
    
    def get_layout(self, layout_type: str) -> SlideLayout:
        """Get layout configuration by type"""
        if layout_type not in self.layouts:
            raise ValueError(f"Unknown layout type: {layout_type}. Available: {list(self.layouts.keys())}")
        return self.layouts[layout_type]
    
    def apply_title_formatting(self, title_shape):
        """Apply theme formatting to title text (without clearing content)"""
        if not title_shape.has_text_frame:
            return
            
        text_frame = title_shape.text_frame
        
        # Apply formatting to existing paragraphs
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Set font properties
            font = paragraph.font
            font.name = self.fonts.title_font
            font.size = Pt(self.fonts.title_size)
            font.color.rgb = self.colors.get_rgb_color('text_dark')
            font.bold = True
    
    def apply_subtitle_formatting(self, subtitle_shape):
        """Apply theme formatting to subtitle text (without clearing content)"""
        if not subtitle_shape.has_text_frame:
            return
            
        text_frame = subtitle_shape.text_frame
        
        # Apply formatting to existing paragraphs
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Set font properties
            font = paragraph.font
            font.name = self.fonts.body_font
            font.size = Pt(self.fonts.subtitle_size)
            font.color.rgb = self.colors.get_rgb_color('text_dark')
    
    def apply_body_formatting(self, body_shape):
        """Apply theme formatting to body text (without clearing content)"""
        if not body_shape.has_text_frame:
            return
            
        text_frame = body_shape.text_frame
        
        # Set text frame properties for optimal space usage
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)  
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        # Ensure text stays within bounds
        text_frame.word_wrap = True
        
        # Apply formatting to existing paragraphs
        for paragraph in text_frame.paragraphs:
            font = paragraph.font
            font.name = self.fonts.body_font
            font.size = Pt(self.fonts.body_size)
            font.color.rgb = self.colors.get_rgb_color('text_dark')
    
    def apply_bullet_formatting(self, paragraph, level: int = 0):
        """Apply bullet point formatting to a paragraph"""
        # Set bullet level (0-based)
        paragraph.level = level
        
        # Set font properties - optimized for content density
        font = paragraph.font
        font.name = self.fonts.body_font
        font.size = Pt(13)  # Slightly smaller for better fit
        font.color.rgb = self.colors.get_rgb_color('text_dark')
        
        # Compact spacing to fit more content
        paragraph.space_after = Pt(4)  # Reduced spacing
        paragraph.space_before = Pt(0)  # No space before bullets
        
        # Allow text wrapping
        paragraph.word_wrap = True


# Predefined themes
class ModernTheme(Theme):
    """Modern blue theme"""
    def __init__(self):
        colors = ThemeColors(
            primary=(68, 114, 196),
            secondary=(112, 173, 71),
            accent=(255, 192, 0)
        )
        super().__init__("Modern", colors)


class CorporateTheme(Theme):
    """Professional corporate theme"""
    def __init__(self):
        colors = ThemeColors(
            primary=(54, 96, 146),
            secondary=(149, 179, 215),
            accent=(180, 198, 231),
            text_dark=(47, 47, 47)
        )
        fonts = ThemeFonts(
            title_font="Segoe UI",
            body_font="Segoe UI"
        )
        super().__init__("Corporate", colors, fonts)


class MinimalistTheme(Theme):
    """Clean minimalist theme"""
    def __init__(self):
        colors = ThemeColors(
            primary=(0, 0, 0),
            secondary=(128, 128, 128),
            accent=(255, 87, 51),
            text_dark=(51, 51, 51)
        )
        fonts = ThemeFonts(
            title_font="Arial",
            body_font="Arial",
            title_size=36,
            subtitle_size=24,
            body_size=16
        )
        super().__init__("Minimalist", colors, fonts)


# Theme factory
AVAILABLE_THEMES = {
    'modern': ModernTheme,
    'corporate': CorporateTheme,
    'minimalist': MinimalistTheme
}


def get_theme(theme_name: str = 'modern') -> Theme:
    """
    Get a theme by name
    
    Args:
        theme_name: Name of the theme ('modern', 'corporate', 'minimalist')
    
    Returns:
        Theme instance
    """
    theme_name = theme_name.lower()
    if theme_name not in AVAILABLE_THEMES:
        raise ValueError(f"Unknown theme: {theme_name}. Available: {list(AVAILABLE_THEMES.keys())}")
    
    return AVAILABLE_THEMES[theme_name]()


# Additional simple Theme dataclass as requested
@dataclass
class SimpleTheme:
    """Simple theme dataclass with basic visual properties"""
    font_name: str
    primary_color: str
    secondary_color: str


# Pre-configured simple theme instances
DEFAULT_THEME = SimpleTheme(
    font_name='Calibri',
    primary_color='0x003366',
    secondary_color='0x336699'
)

DARK_THEME = SimpleTheme(
    font_name='Arial',
    primary_color='0xFFFFFF',
    secondary_color='0xCCCCCC'
)
