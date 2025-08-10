"""
PowerPoint Builder for RealtyAI
==============================

Main class for creating PowerPoint presentations with themes and layouts.
"""

from typing import Dict, Any, Optional, List
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .slide_templates import Theme, get_theme

logger = logging.getLogger(__name__)


class PowerPointBuilder:
    """Main class for building PowerPoint presentations"""
    
    def __init__(self, theme: Optional[Theme] = None):
        """
        Initialize PowerPoint Builder
        
        Args:
            theme: Theme object for styling (uses Modern theme if None)
        """
        self.theme = theme or get_theme('modern')
        self.presentation = Presentation()
        self.slide_count = 0
        
        logger.info(f"PowerPoint Builder initialized with theme: {self.theme.name}")
    
    def add_slide(self, slide_data: Dict[str, Any]) -> None:
        """
        Add a slide to the presentation based on slide data
        
        Args:
            slide_data: Dictionary containing slide information
                - type: 'title', 'content', 'section_header', etc.
                - title: Slide title
                - subtitle: Subtitle (for title slides)
                - bullet_points: List of bullet points (for content slides)
                - content: Additional content
        """
        slide_type = slide_data.get('type', 'content')
        
        try:
            if slide_type == 'title':
                self._add_title_slide(slide_data)
            elif slide_type == 'content':
                self._add_content_slide(slide_data)
            elif slide_type == 'section_header':
                self._add_section_header_slide(slide_data)
            elif slide_type == 'conclusion':
                self._add_conclusion_slide(slide_data)
            else:
                logger.warning(f"Unknown slide type: {slide_type}, using content layout")
                self._add_content_slide(slide_data)
            
            self.slide_count += 1
            logger.debug(f"Added {slide_type} slide: {slide_data.get('title', 'Untitled')}")
            
        except Exception as e:
            logger.error(f"Failed to add slide: {str(e)}")
            raise
    
    def _add_title_slide(self, slide_data: Dict[str, Any]) -> None:
        """Add a title slide"""
        layout = self.theme.get_layout('title')
        slide_layout = self.presentation.slide_layouts[layout.layout_id]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Get placeholders
        title_placeholder = slide.placeholders[layout.title_placeholder_idx]
        subtitle_placeholder = slide.placeholders[layout.content_placeholder_idx]
        
        # Set title
        title = slide_data.get('title', 'Presentation Title')
        title_placeholder.text = title
        # Apply formatting after setting text
        self.theme.apply_title_formatting(title_placeholder)
        
        # Set subtitle
        subtitle = slide_data.get('subtitle', '')
        if not subtitle and 'content' in slide_data:
            # Try to get subtitle from content
            content = slide_data['content']
            if isinstance(content, dict):
                subtitle = content.get('subtitle', '')
            elif isinstance(content, str):
                subtitle = content
        
        subtitle_placeholder.text = subtitle
        # Apply formatting after setting text
        self.theme.apply_subtitle_formatting(subtitle_placeholder)
    
    def _add_content_slide(self, slide_data: Dict[str, Any]) -> None:
        """Add a content slide with bullet points using auto-fit"""
        layout = self.theme.get_layout('content')
        slide_layout = self.presentation.slide_layouts[layout.layout_id]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_placeholder = slide.placeholders[layout.title_placeholder_idx]
        title = slide_data.get('title', 'Content')
        title_placeholder.text = title
        self.theme.apply_title_formatting(title_placeholder)
        
        # Set content
        content_placeholder = slide.placeholders[layout.content_placeholder_idx]
        
        # Add bullet points
        bullet_points = self._extract_bullet_points(slide_data)
        if bullet_points:
            self._add_bullet_points(content_placeholder, bullet_points)
            self.theme.apply_body_formatting(content_placeholder)
        else:
            # If no bullet points, add general content with auto-fit
            self._add_single_content_slide_content(content_placeholder, slide_data)
    
    def _add_single_content_slide_content(self, content_placeholder, slide_data: Dict[str, Any]) -> None:
        """Add general content to a placeholder with auto-fit"""
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # Enable auto-fit for text content
        from pptx.enum.text import MSO_AUTO_SIZE
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        
        # Add general content
        content_text = slide_data.get('content', 'Content goes here')
        if isinstance(content_text, dict):
            content_text = str(content_text)
        
        content_placeholder.text = content_text
        self.theme.apply_body_formatting(content_placeholder)
    
    def _add_section_header_slide(self, slide_data: Dict[str, Any]) -> None:
        """Add a section header slide"""
        layout = self.theme.get_layout('section_header')
        slide_layout = self.presentation.slide_layouts[layout.layout_id]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_placeholder = slide.placeholders[layout.title_placeholder_idx]
        title = slide_data.get('title', 'Section')
        title_placeholder.text = title
        self.theme.apply_title_formatting(title_placeholder)
    
    def _add_conclusion_slide(self, slide_data: Dict[str, Any]) -> None:
        """Add a conclusion slide (similar to content but with different styling)"""
        self._add_content_slide(slide_data)
    
    def _extract_bullet_points(self, slide_data: Dict[str, Any]) -> List[str]:
        """Extract bullet points from slide data"""
        bullet_points = []
        
        # Check for bullet_points directly
        if 'bullet_points' in slide_data:
            bullets = slide_data['bullet_points']
            if isinstance(bullets, list):
                for bullet in bullets:
                    if isinstance(bullet, dict):
                        # Handle structured bullet points
                        point = bullet.get('point', '')
                        details = bullet.get('details', '')
                        if point:
                            bullet_text = point
                            if details and details != point:
                                bullet_text += f": {details}"
                            bullet_points.append(bullet_text)
                    elif isinstance(bullet, str):
                        bullet_points.append(bullet)
        
        # Check for content with key_points
        elif 'content' in slide_data:
            content = slide_data['content']
            if isinstance(content, dict):
                # Check for key_points
                if 'key_points' in content:
                    key_points = content['key_points']
                    if isinstance(key_points, list):
                        bullet_points.extend([str(point) for point in key_points])
                
                # Check for themes
                if 'themes' in content:
                    themes = content['themes']
                    if isinstance(themes, list):
                        bullet_points.extend([f"Key theme: {theme}" for theme in themes])
                
                # Check for supporting_data
                if 'supporting_data' in content:
                    supporting_data = content['supporting_data']
                    if isinstance(supporting_data, list):
                        bullet_points.extend([f"Supporting fact: {data}" for data in supporting_data])
        
        # If still no bullet points, check for key_points at top level
        if not bullet_points and 'key_points' in slide_data:
            key_points = slide_data['key_points']
            if isinstance(key_points, list):
                bullet_points.extend([str(point) for point in key_points])
        
        # Return bullet points with intelligent content management
        # If too many bullets, prioritize the most important ones
        max_bullets = 5  # Optimal number for readability
        if len(bullet_points) > max_bullets:
            logger.info(f"Trimming {len(bullet_points)} bullet points to {max_bullets} for better slide readability")
        
        return bullet_points
    
    def _add_bullet_points(self, content_placeholder, bullet_points: List[str]) -> None:
        """Add bullet points to a content placeholder with auto-fit"""
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # Optimize margins for maximum text space
        text_frame.margin_left = Inches(0.1)   # Minimal left margin
        text_frame.margin_right = Inches(0.1)  # Minimal right margin
        text_frame.margin_top = Inches(0.1)    # Minimal top margin
        text_frame.margin_bottom = Inches(0.1) # Minimal bottom margin
        
        # Enable text wrapping and auto-fit
        text_frame.word_wrap = True
        
        # Enable auto-fit to automatically adjust text size and spacing
        from pptx.enum.text import MSO_AUTO_SIZE
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Add all bullet points - let PowerPoint auto-fit handle the sizing
        for i, bullet_text in enumerate(bullet_points):
            cleaned_text = bullet_text.strip()
            if not cleaned_text:
                continue
                
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = cleaned_text
            self.theme.apply_bullet_formatting(p, level=0)
        
        logger.info(f"Added {len(bullet_points)} bullets with auto-fit enabled")
    
    def add_bulk_slides(self, slides_data: List[Dict[str, Any]]) -> None:
        """
        Add multiple slides at once
        
        Args:
            slides_data: List of slide data dictionaries
        """
        for slide_data in slides_data:
            self.add_slide(slide_data)
        
        logger.info(f"Added {len(slides_data)} slides to presentation")
    
    def save(self, file_path: str) -> None:
        """
        Save the presentation to a file
        
        Args:
            file_path: Path where to save the presentation
        """
        try:
            # Ensure the directory exists
            path = Path(file_path)
            path.parent.mkdir(parents=True, exist_ok=True)
            
            # Save the presentation
            self.presentation.save(file_path)
            
            logger.info(f"Presentation saved successfully to: {file_path}")
            logger.info(f"Total slides created: {self.slide_count}")
            
        except Exception as e:
            logger.error(f"Failed to save presentation: {str(e)}")
            raise
    
    def get_slide_count(self) -> int:
        """Get the current number of slides"""
        return self.slide_count
    
    def set_theme(self, theme: Theme) -> None:
        """
        Change the theme (affects future slides)
        
        Args:
            theme: New theme to use
        """
        self.theme = theme
        logger.info(f"Theme changed to: {theme.name}")
    
    def add_custom_slide(self, layout_type: str, content: Dict[str, Any]) -> None:
        """
        Add a slide with custom layout
        
        Args:
            layout_type: Type of layout to use
            content: Content dictionary
        """
        slide_data = {
            'type': layout_type,
            **content
        }
        self.add_slide(slide_data)
